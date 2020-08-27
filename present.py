from pptx import Presentation
import pandas as pd
import os
import subprocess
from datetime import datetime


class excelPowerpoint:
    def __init__(self):
        pass
    #Function to extract data from inventory spreadsheet
    def find_data(self, location, inventory_number):
        inventories = ['no_data']
        for i in range(1,5):
            df = pd.read_excel('Inventory Workbook.xlsx', sheet_name=f'Inventory{i}')
            inventories.append(df)

        for i in range(len(inventories[inventory_number]['Location (Suite, Rack, Shelf)'])):
            if inventories[inventory_number]['Site'][i] == location:
                return str(inventories[inventory_number]['Location (Suite, Rack, Shelf)'][i])

        raise ValueError(f'Location "{location}" not on spreadsheet')

    def findshapenumber(self, searchvalue):
        prs = Presentation('template.pptx')
        slide = prs.slides[0]
        for i in range(len(slide.shapes)):
            try:
                if slide.shapes[i].text == searchvalue:
                    return i
            except AttributeError:
                pass
        raise ValueError('Shape not found')

#User enter Primary and Secondary Sites
print("Welcome to Design Automation")
Firstsite=input("Please Enter Primary Site Name: ")
Secondsite=input("Please Enter Second Site Name: ")

excel_powerpoint = excelPowerpoint()
#Attributes to get from excel
FirstsitePosA = excel_powerpoint.find_data(Firstsite, 1)
FirstsitePosB = excel_powerpoint.find_data(Firstsite, 2)
SecondsitePosA = excel_powerpoint.find_data(Secondsite, 4)

#Editing presentation template with excel attributes:
prs = Presentation('template.pptx')
slide = prs.slides[0]
slide.shapes[excel_powerpoint.findshapenumber('LocationLabel2')].text = FirstsitePosB
slide.shapes[excel_powerpoint.findshapenumber('LocationLabel1')].text = FirstsitePosA
slide.shapes[excel_powerpoint.findshapenumber('LocationLabel3')].text = SecondsitePosA
slide.shapes[excel_powerpoint.findshapenumber('LocationLabel4')].text = Firstsite
slide.shapes[excel_powerpoint.findshapenumber('LocationLabel5')].text = Secondsite

#Saving the edited presentation
file_name = f'{Firstsite}_{Secondsite}_{datetime.now()}.pptx'
prs.save(file_name)
opener = 'open'
subprocess.call([opener, file_name])




